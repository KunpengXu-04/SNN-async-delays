"""
make_presentation_flowcharts.py
Build two bridging slides that visually connect Steps 1-4 and the three new
topologies, using the same visual language as the existing supervisor deck.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from make_presentation import (
    W,
    H,
    C_NAVY,
    C_IC,
    C_ORANGE,
    C_WHITE,
    C_BLACK,
    C_ALT,
    CONT_LEFT,
    CONT_TOP,
    CONT_W,
    add_textbox,
    add_rect,
    add_title_bar,
    add_bullet_box,
    add_callout_box,
    new_slide,
)


ROOT = os.path.dirname(__file__)
OUT_PATH = os.path.join(ROOT, "presentation_topology_flowcharts.pptx")
ROADMAP_IMG = os.path.join(ROOT, "topology 关系图.png")


def _set_line(shape, color, width_pt=1.5):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _rounded_box(slide, left, top, width, height, text, fill_rgb, line_rgb, font_size=15, bold=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    _set_line(shape, line_rgb, 1.5)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = C_NAVY
    return shape


def _circle_box(slide, left, top, diameter, text, fill_rgb, line_rgb, font_size=15, bold=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, left, top, diameter, diameter
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    _set_line(shape, line_rgb, 1.5)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = C_NAVY
    return shape


def _arrow(slide, x1, y1, x2, y2, color=C_ORANGE):
    left = min(x1, x2)
    top = min(y1, y2) - Inches(0.08)
    width = max(Inches(0.35), abs(x2 - x1))
    height = Inches(0.16)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide1_roadmap(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "How Step 4 Connects to Steps 1-3")

    add_callout_box(
        slide,
        "These topologies are not a brand-new branch of the project: they are a controlled refactor of the same question. "
        "Each one isolates a different bottleneck that earlier steps mixed together.",
        CONT_LEFT,
        CONT_TOP,
        CONT_W,
        Inches(0.9),
        bg_color=C_ALT,
        border_color=C_IC,
        font_size=14,
    )

    if os.path.exists(ROADMAP_IMG):
        slide.shapes.add_picture(
            ROADMAP_IMG,
            Inches(0.55),
            Inches(2.0),
            width=Inches(6.9),
            height=Inches(4.45),
        )
    else:
        add_callout_box(
            slide,
            "Roadmap image missing: topology 关系图.png",
            Inches(0.55),
            Inches(2.0),
            Inches(6.9),
            Inches(4.45),
            bg_color=RGBColor(0xFE, 0xF2, 0xF2),
            border_color=C_ORANGE,
            font_size=16,
        )

    bullets = [
        ("Step 1: one query, one op, one output. Baseline solvability.", 0),
        ("Step 2: many queries, one op. Plan A / C / D change how queries share time and channels.", 0),
        ("Step 3: many queries, many ops, many outputs. Mixed-task setting without output compression.", 0),
        ("Topology 1: removes input collision, keeps multiple outputs.", 0),
        ("Topology 2: keeps shared-channel input, compresses K answers into one final decision.", 0),
        ("Topology 3: mixes different ops and also compresses to one final decision.", 0),
    ]
    add_bullet_box(
        slide,
        bullets,
        Inches(7.8),
        Inches(2.0),
        Inches(4.7),
        Inches(3.55),
        font_size=15,
    )

    add_callout_box(
        slide,
        "Reading guide: Topology 1 asks whether delays are still needed when the same input is broadcast once. "
        "Topology 2 asks what changes when K per-query answers must be fused into one scalar output. "
        "Topology 3 stacks both difficulties together.",
        Inches(7.8),
        Inches(5.75),
        Inches(4.7),
        Inches(0.95),
        bg_color=RGBColor(0xE8, 0xF0, 0xFE),
        border_color=C_IC,
        font_size=13,
    )


def _panel_label(slide, top, label):
    add_textbox(
        slide,
        label,
        Inches(0.55),
        top,
        Inches(1.1),
        Inches(0.3),
        font_size=17,
        bold=True,
        color=C_NAVY,
        align=PP_ALIGN.LEFT,
    )


def _shared_snn_box(slide, left, top, width, height, subtitle):
    box = _rounded_box(
        slide,
        left,
        top,
        width,
        height,
        f"Shared hidden SNN\n{subtitle}",
        RGBColor(0xE8, 0xF0, 0xFE),
        C_IC,
        font_size=15,
        bold=True,
    )
    return box


def slide2_structures(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Topology Structures at a Glance")

    row_top = [Inches(1.45), Inches(3.08), Inches(4.72)]
    row_h = Inches(1.25)

    for top in row_top:
        panel = add_rect(slide, Inches(0.45), top, Inches(12.4), row_h, C_WHITE, line_color=RGBColor(0xD8, 0xE2, 0xEE), line_width=1)
        panel.fill.transparency = 0

    _panel_label(slide, Inches(1.72), "T1")
    _panel_label(slide, Inches(3.35), "T2")
    _panel_label(slide, Inches(4.99), "T3")

    q1 = _rounded_box(slide, Inches(1.2), Inches(1.58), Inches(2.0), Inches(0.65), "One query\n(A, B)", C_ALT, C_IC, bold=True)
    snn1 = _shared_snn_box(slide, Inches(4.0), Inches(1.48), Inches(2.25), Inches(0.88), "No input collision")
    head1 = _rounded_box(slide, Inches(7.0), Inches(1.45), Inches(1.25), Inches(0.45), "AND", RGBColor(0xF8, 0xE6, 0xD5), C_ORANGE, bold=True)
    head2 = _rounded_box(slide, Inches(8.4), Inches(1.45), Inches(1.25), Inches(0.45), "OR", RGBColor(0xF8, 0xE6, 0xD5), C_ORANGE, bold=True)
    head3 = _rounded_box(slide, Inches(9.8), Inches(1.45), Inches(1.25), Inches(0.45), "NAND", RGBColor(0xF8, 0xE6, 0xD5), C_ORANGE, bold=True)
    head4 = _rounded_box(slide, Inches(11.2), Inches(1.45), Inches(1.0), Inches(0.45), "NOR", RGBColor(0xF8, 0xE6, 0xD5), C_ORANGE, bold=True)
    add_textbox(slide, "K independent readout heads", Inches(7.0), Inches(2.0), Inches(5.0), Inches(0.28), font_size=12, color=C_NAVY, align=PP_ALIGN.CENTER)
    _arrow(slide, q1.left + q1.width, q1.top + q1.height / 2, snn1.left, snn1.top + snn1.height / 2)
    _arrow(slide, snn1.left + snn1.width, snn1.top + snn1.height / 2, head1.left, head1.top + head1.height / 2)

    q2a = _rounded_box(slide, Inches(1.2), Inches(3.19), Inches(1.25), Inches(0.42), "q1", C_ALT, C_IC, bold=True)
    q2b = _rounded_box(slide, Inches(1.2), Inches(3.66), Inches(1.25), Inches(0.42), "q2", C_ALT, C_IC, bold=True)
    q2c = _rounded_box(slide, Inches(1.2), Inches(4.13), Inches(1.25), Inches(0.42), "q3", C_ALT, C_IC, bold=True)
    add_textbox(slide, "Sequential sub-windows on 2 shared channels", Inches(2.65), Inches(3.47), Inches(2.1), Inches(0.35), font_size=12, color=C_NAVY, align=PP_ALIGN.CENTER)
    snn2 = _shared_snn_box(slide, Inches(5.0), Inches(3.12), Inches(2.35), Inches(0.98), "Delay-based routing")
    agg2 = _circle_box(slide, Inches(9.45), Inches(3.19), Inches(0.95), "1\nagg", RGBColor(0xE7, 0xF4, 0xEA), RGBColor(0x1A, 0x7A, 0x3C), font_size=14, bold=True)
    add_textbox(slide, "majority or AND over K per-query answers", Inches(10.6), Inches(3.37), Inches(2.0), Inches(0.45), font_size=12, color=C_NAVY)
    _arrow(slide, Inches(2.55), Inches(3.82), snn2.left, snn2.top + snn2.height / 2)
    _arrow(slide, snn2.left + snn2.width, snn2.top + snn2.height / 2, agg2.left, agg2.top + agg2.height / 2)

    q3a = _rounded_box(slide, Inches(1.2), Inches(4.82), Inches(1.45), Inches(0.42), "q1 + op1", C_ALT, C_IC, bold=True)
    q3b = _rounded_box(slide, Inches(1.2), Inches(5.29), Inches(1.45), Inches(0.42), "q2 + op2", C_ALT, C_IC, bold=True)
    q3c = _rounded_box(slide, Inches(1.2), Inches(5.76), Inches(1.45), Inches(0.42), "q3 + op3", C_ALT, C_IC, bold=True)
    add_textbox(slide, "Shared channels + mixed op identity", Inches(2.85), Inches(5.11), Inches(2.15), Inches(0.35), font_size=12, color=C_NAVY, align=PP_ALIGN.CENTER)
    snn3 = _shared_snn_box(slide, Inches(5.0), Inches(4.75), Inches(2.35), Inches(0.98), "Routing + op separation")
    agg3 = _circle_box(slide, Inches(9.45), Inches(4.83), Inches(0.95), "1\nagg", RGBColor(0xE7, 0xF4, 0xEA), RGBColor(0x1A, 0x7A, 0x3C), font_size=14, bold=True)
    add_textbox(slide, "one final decision after mixed-op processing", Inches(10.6), Inches(5.0), Inches(2.0), Inches(0.45), font_size=12, color=C_NAVY)
    _arrow(slide, Inches(2.75), Inches(5.43), snn3.left, snn3.top + snn3.height / 2)
    _arrow(slide, snn3.left + snn3.width, snn3.top + snn3.height / 2, agg3.left, agg3.top + agg3.height / 2)

    add_callout_box(
        slide,
        "Aggregate output means the network no longer emits K separate answers. Instead, the readout must compress the whole trial into one binary verdict, "
        "such as whether the majority of queries are correct-1 or whether all queried NAND results are 1.",
        Inches(0.7),
        Inches(6.35),
        Inches(12.0),
        Inches(0.7),
        bg_color=RGBColor(0xF9, 0xF6, 0xEF),
        border_color=C_ORANGE,
        font_size=13,
    )


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide1_roadmap(prs)
    slide2_structures(prs)
    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    build_deck()
