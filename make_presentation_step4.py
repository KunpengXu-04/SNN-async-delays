"""
make_presentation_step4.py
Generate supervisor update PPTX covering work since the last meeting (Step 4:
Pattern Taxonomy — Topologies 1-3, plus Section 27 depth ablation). Reuses the
style/helpers from make_presentation.py and embeds real result figures
(spike-flow plots, diagnostic panels, energy chart) instead of only tables.
Run: python make_presentation_step4.py
"""

import os
import sys
import tempfile

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

sys.path.insert(0, os.path.dirname(__file__))
from make_presentation import (
    add_textbox, add_rect, add_title_bar, new_slide, make_table,
    add_bullet_box, add_callout_box,
    C_NAVY, C_IC, C_ORANGE, C_WHITE, C_BLACK, C_ALT, C_GREEN, C_RED,
    CONT_TOP, CONT_LEFT, CONT_W, CONT_H, W, H,
)

ROOT = os.path.dirname(__file__)
SNN = os.path.join(ROOT, "snn_async_delays")
OUT_PATH = os.path.join(ROOT, "presentation_step4_update.pptx")

# Real result figures (already rendered by the experiment pipeline — embed as-is)
FIG_PLAND_K3_WAD = os.path.join(SNN, "paper", "figures", "fig_planD_k3_mlp_flow.png")
FIG_PLAND_K3_D0 = os.path.join(SNN, "paper", "figures", "fig_ablation_mlp_d0_flow.png")
FIG_DIAGNOSTIC_T3 = os.path.join(
    SNN, "runs", "many_many_one_out_(step4)", "wad_mlp_K3_seed42", "plots", "diagnostic_panel.png")
FIG_ENERGY_CHART = os.path.join(SNN, "docs", "computation_vs_energy_v2.png")


# ---------------------------------------------------------------------------
# Generic grouped bar-chart helper (used by T1 / T2 / T3 / Section 27 slides)
# ---------------------------------------------------------------------------
def make_bar_chart(tmp_path, k_labels, series, title, ylabel,
                    threshold=None, threshold_label="90% threshold",
                    colors=None, figsize=(6.6, 4.0)):
    """series: dict {label: [values per k_labels]} (values may contain None)."""
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    n_groups = len(k_labels)
    n_series = len(series)
    bar_w = 0.8 / n_series
    x = np.arange(n_groups)

    default_colors = ["#003E74", "#E8600A", "#888888", "#1A7A3C"]
    colors = colors or default_colors

    for i, (label, vals) in enumerate(series.items()):
        y = [v if v is not None else 0 for v in vals]
        offset = (i - (n_series - 1) / 2) * bar_w
        ax.bar(x + offset, y, width=bar_w, label=label,
               color=colors[i % len(colors)])

    if threshold is not None:
        ax.axhline(threshold, color="#C0142C", linewidth=1.5, linestyle="-.",
                   label=threshold_label)

    ax.set_xticks(x)
    ax.set_xticklabels(k_labels, fontsize=11)
    ax.set_ylabel(ylabel, fontsize=12)
    ax.set_title(title, fontsize=13, fontweight="bold", color="#1F3864")
    ax.legend(fontsize=9, loc="lower right", framealpha=0.9)
    ax.grid(True, axis="y", alpha=0.3, linestyle="--")
    for spine in ax.spines.values():
        spine.set_edgecolor("#CCCCCC")

    plt.tight_layout()
    plt.savefig(tmp_path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def add_chart_picture(slide, chart_name, left, top, width, height, build_fn):
    tmp = os.path.join(tempfile.gettempdir(), f"tmp_chart_{chart_name}.png")
    build_fn(tmp)
    slide.shapes.add_picture(tmp, left, top, width, height)
    try:
        os.remove(tmp)
    except Exception:
        pass


def add_image(slide, path, left, top, width=None, height=None):
    """Add an existing image file, preserving aspect ratio if only one dim given."""
    slide.shapes.add_picture(path, left, top, width=width, height=height)


# ---------------------------------------------------------------------------
def slide1_title(prs):
    slide = new_slide(prs)
    add_rect(slide, Inches(0), Inches(0), W, Inches(4.2), C_NAVY)
    add_rect(slide, Inches(0), Inches(4.2), W, Inches(3.3), C_WHITE)
    add_rect(slide, Inches(0), Inches(4.1), W, Inches(0.12), C_ORANGE)

    add_textbox(slide,
                "Step 4: Pattern Taxonomy —\nWhen Are Delays Actually Necessary?",
                Inches(0.7), Inches(0.7), Inches(11.9), Inches(2.8),
                font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                "Update since last meeting  ·  Three new output-structure topologies "
                "+ a depth-ablation correction",
                Inches(0.7), Inches(3.35), Inches(11.9), Inches(0.7),
                font_size=18, color=C_ALT, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                "SNN Async Delays Project  ·  June 2026",
                Inches(0.7), Inches(4.45), Inches(8), Inches(0.55),
                font_size=15, color=C_NAVY, align=PP_ALIGN.LEFT)
    add_textbox(slide,
                "Imperial College London",
                Inches(9.0), Inches(4.45), Inches(4), Inches(0.55),
                font_size=14, color=C_IC, align=PP_ALIGN.RIGHT)
    return slide


def slide2_recap(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Recap: Where We Left Off")

    add_callout_box(slide,
        "Plan D (shared-channel, sequential queries) established the central claim: "
        "trainable delays are structurally necessary when K queries share both input "
        "channels and a single readout window. d=0 plateaus at ~77% regardless of "
        "readout type; MLP readout pushes Max K@90% from 2 to 3.",
        CONT_LEFT, CONT_TOP, CONT_W, Inches(1.2),
        bg_color=RGBColor(0xE8, 0xF0, 0xFE), border_color=C_IC, font_size=15)

    add_textbox(slide,
                "Open question raised at last meeting: is this necessity a general property of "
                "shared-channel SNNs, or an artifact of this one input/output structure?",
                CONT_LEFT, Inches(2.35), CONT_W, Inches(0.6),
                font_size=15, bold=True, color=C_NAVY)

    bullets = [
        ("Step 4 answers this by varying the I/O structure while holding the core "
         "mechanism (shared channels, sequential sub-windows, LIF + delays) fixed:", 0),
        ("Topology 1 — one-query, many-op:  1 broadcast input  →  K_ops independent outputs", 1),
        ("Topology 2 — many-query, one-out:  K sequential inputs  →  1 aggregate output", 1),
        ("Topology 3 — many-query, many-op, one-out:  combines both (mixed ops + aggregation)", 1),
        ("Plus: a depth ablation revisited the single-op (NAND) ceiling and found a new best result.", 0),
    ]
    add_bullet_box(slide, bullets, CONT_LEFT, Inches(3.1), CONT_W, Inches(3.0), font_size=16)
    return slide


def slide3_taxonomy(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "The Pattern Taxonomy")

    headers = ["Topology", "Input", "Output", "Tests"]
    rows = [
        ["Plan D (done)", "K shared-channel\nsequential", "K independent logits",
         "Baseline: delay-as-router\nnecessity"],
        ["T1: one-query-many-op", "1 broadcast pair", "K_ops independent\nlogits",
         "Does necessity require\ninput-side collision?"],
        ["T2: many-query-one-out", "K shared-channel\nsequential", "1 aggregate logit\n(majority / AND)",
         "Does necessity survive\noutput-side compression?"],
        ["T3: many-query-many-op\n-one-out", "K shared-channel,\nmixed ops", "1 aggregate logit",
         "Do routing + op-mix +\naggregation compound?"],
    ]
    hl = {(0, c): RGBColor(0xEE, 0xF2, 0xF7) for c in range(4)}
    make_table(slide, headers, rows,
               CONT_LEFT, CONT_TOP, CONT_W, Inches(3.6),
               col_widths=[2.4, 2.6, 2.6, 3.0],
               font_size=13, highlight_cells=hl)

    add_callout_box(slide,
        "Design logic: Plan D conflated two things — input-side temporal collision and "
        "output-side multiplicity. T1 isolates the output side (no input collision); "
        "T2 isolates the compression side (input collision, but 1 output); T3 stacks all three.",
        CONT_LEFT, Inches(5.1), CONT_W, Inches(1.1),
        bg_color=C_ALT, border_color=C_IC, font_size=14)
    return slide


def slide4_t1(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Topology 1: One-Query, Many-Op — Delays Are NOT Necessary")

    add_callout_box(slide,
        "Design: single shared (A,B) pair in ONE input window  →  K_ops readout heads, "
        "each learning a different op (AND/OR/NAND/NOR) of the same input. No input-side "
        "collision — each head already has its own readout weights.",
        CONT_LEFT, CONT_TOP, CONT_W, Inches(0.8),
        bg_color=C_ALT, border_color=C_IC, font_size=13)

    def build(tmp):
        make_bar_chart(tmp,
            ["K_ops=2", "K_ops=3", "K_ops=4"],
            {
                "wad_linear": [96.1, 95.9, 96.1],
                "wad_mlp":    [94.4, 94.2, 94.2],
                "d0_mlp":     [80.8, 81.2, 79.3],
                "d0_linear":  [77.6, 78.7, 78.1],
            },
            "Topology 1: Accuracy vs K_ops", "Accuracy (%)",
            threshold=90)

    add_chart_picture(slide, "t1", Inches(0.5), Inches(1.95), Inches(7.0), Inches(4.25), build)

    bullets = [
        ("wad-vs-d0 gap is flat (+14–19pp) across K_ops — does NOT grow, unlike Plan D (+16%→+34%).", 0),
        ("Linear readout beats MLP here (reversal of Plan D): K_ops heads already separate; "
         "MLP's nonlinear mixing adds inter-head competition instead of helping.", 0),
        ("No capacity ceiling found within K_ops ≤ 4 — wad_linear stays ≥95% throughout.", 0),
        ("d=0 collapses to the label prior (OR/NOR plateau exactly at 75.6% = P(label=1)).", 0),
    ]
    add_bullet_box(slide, bullets, Inches(7.7), Inches(2.1), Inches(5.1), Inches(4.0), font_size=14)
    return slide


def slide5_t2(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Topology 2: Many-Query, One-Out — Delays ARE Necessary")

    add_callout_box(slide,
        "Design: keeps Plan D's input side exactly (K NAND queries, sequential sub-windows, "
        "2 shared channels) but collapses the K-logit readout to ONE aggregate logit "
        "— majority vote, or AND-of-results (cleaner: trivial baseline ≈ 50% balanced acc).",
        CONT_LEFT, CONT_TOP, CONT_W, Inches(0.9),
        bg_color=C_ALT, border_color=C_IC, font_size=13)

    def build(tmp):
        make_bar_chart(tmp,
            ["K=1", "K=3", "K=5"],
            {
                "Majority — wad_mlp":    [93.6, 76.4, 64.8],
                "Majority — d0 (any)":   [70.4, 50.0, 50.0],
                "AND — wad_mlp":         [None, 86.8, 75.0],
                "AND — d0 (any)":        [None, 50.0, 50.0],
            },
            "Topology 2: Balanced Accuracy vs K", "Balanced Accuracy (%)",
            threshold=70, threshold_label="70% BalAcc threshold")

    add_chart_picture(slide, "t2", Inches(0.5), Inches(2.0), Inches(7.0), Inches(4.2), build)

    bullets = [
        ("d=0 balanced accuracy ≈50% at K≥3 for both aggregations — identical structural "
         "failure to Plan D (confirmed via spike-flow plots: zero readout-window activity).", 0),
        ("MLP beats linear here (opposite of T1): aggregating K sequential signals into one "
         "decision needs nonlinear counting/thresholding, not K independent linear sums.", 0),
        ("Max K@BalAcc≥70%: K=3 (majority vote, matches Plan D's ceiling); "
         "K=5 (AND-of-results — conjunctive structure extends the ceiling further).", 0),
    ]
    add_bullet_box(slide, bullets, Inches(7.7), Inches(2.15), Inches(5.1), Inches(4.0), font_size=14)
    return slide


def slide6_t3(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Topology 3: Mixed-Op + Aggregation — Bottlenecks Compound")

    add_callout_box(slide,
        "Design: Topology 2's K→1 majority vote combined with Step 3's mixed-op input "
        "(4 ops, h=100, n_train=16000) — the most demanding composition in the taxonomy.",
        CONT_LEFT, CONT_TOP, CONT_W, Inches(0.65),
        bg_color=C_ALT, border_color=C_IC, font_size=13)

    def build(tmp):
        make_bar_chart(tmp,
            ["K=1", "K=3"],
            {
                "wad_mlp": [97.0, 80.9],
                "d0_mlp":  [65.6, 57.0],
            },
            "Topology 3: Accuracy vs K", "Accuracy (%)",
            threshold=90, figsize=(5.6, 3.6))

    add_chart_picture(slide, "t3", CONT_LEFT, Inches(1.55), Inches(5.6), Inches(3.4), build)

    add_textbox(slide,
                "Delay gap: +31.4pp (K=1), +23.9pp (K=3)  |  Max K@90%: wad_mlp = 1, d0_mlp = 0",
                CONT_LEFT, Inches(5.05), Inches(5.6), Inches(0.5),
                font_size=13, bold=True, color=C_NAVY)

    # cross-topology comparison (kept as table — compact, categorical, 4 rows)
    headers2 = ["Topology", "Delay needed?", "Max K@90% (wad)", "d0 failure mode"]
    rows2 = [
        ["Plan D",                  "Yes — structural", "3 (MLP, h=50)", "Weight-symmetry, ~77%"],
        ["T1: one-query-many-op",   "No — marginal",     "4 (all pass)",  "Slight alignment loss"],
        ["T2: many-query-one-out",  "Yes — structural",  "K=3 (BalAcc)",  "Label-prior collapse"],
        ["T3: many-many-one-out",   "Yes — structural",  "1",             "Label prior + op-mix"],
    ]
    hl2 = {(1, c): RGBColor(0xFA, 0xE8, 0xE8) for c in range(4)}
    hl2.update({(3, c): RGBColor(0xD6, 0xE8, 0xD6) for c in range(4)})
    make_table(slide, headers2, rows2,
               Inches(6.4), Inches(1.55), Inches(6.4), Inches(2.7),
               col_widths=[2.0, 1.6, 1.6, 2.2],
               font_size=11, highlight_cells=hl2)

    add_textbox(slide,
                "Max K@90%=1 sits below Direction C's K=3 (mixed-op only, h=100) and T2's K=3 "
                "(aggregation only, h=50) — routing + op-mix + aggregation compound multiplicatively, not additively.",
                CONT_LEFT, Inches(5.65), CONT_W, Inches(0.9),
                font_size=13, color=RGBColor(0x44, 0x44, 0x44))
    return slide


def slide7_section27(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Section 27: Depth Revisited — New Best Result for Single-Op NAND")

    add_callout_box(slide,
        "Earlier (Sections 13-14) we concluded a 2nd hidden layer doesn't help — but that "
        "test split a fixed 50-neuron budget into 25+25, halving the readout's input "
        "dimension. This experiment removes that confound: L2-h50h50 uses two FULL "
        "50-neuron layers (readout_in=50, matching L1-h50) — 100 neurons total, no budget cut.",
        CONT_LEFT, CONT_TOP, CONT_W, Inches(1.1),
        bg_color=C_ALT, border_color=C_IC, font_size=13)

    def build(tmp):
        make_bar_chart(tmp,
            ["L1-h50\nlinear", "L1-h50\nMLP", "L2-h50h50\nlinear", "L2-h50h50\nMLP"],
            {"Max K@90%": [2, 3, 4, 3]},
            "Max K@90% — Single-Op NAND", "Max K @ 90% accuracy",
            colors=["#003E74"], figsize=(6.0, 4.0))

    add_chart_picture(slide, "sec27", Inches(0.5), Inches(2.25), Inches(6.6), Inches(4.0), build)

    bullets = [
        ("L2-h50h50-linear reaches Max K@90%=4 — the best single-op NAND result in the "
         "whole project, beating the previous best (L1-h50-MLP, K=3).", 0),
        ("Confirms the earlier \"depth hurts\" finding was an artifact of the readout-input "
         "cut (50→25), not of depth itself: more neurons + matched readout_in together help.", 0),
        ("First reversal of linear<MLP at K=4 (91.1% vs 89.7%): the 2nd spiking layer's own "
         "LIF nonlinearity already supplies the separation MLP used to add on top of L1.", 0),
        ("Open caveat: no L1-h100 single-op control was run, so depth vs. raw neuron count "
         "(100 total either way) isn't fully disentangled — flagged in the paper's Future Work.", 0),
    ]
    add_bullet_box(slide, bullets, Inches(7.4), Inches(2.4), Inches(5.4), Inches(4.0), font_size=13)
    return slide


def slide8_mechanism(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Mechanism: How Delays Route Information in Time")

    add_textbox(slide,
                "Same setup (Plan D, K=3, MLP readout) — only the delay condition differs.",
                CONT_LEFT, CONT_TOP, CONT_W, Inches(0.45),
                font_size=14, bold=True, color=C_NAVY)

    img_h = Inches(5.55)
    img_w = Inches(3.59)  # aspect-preserved (1674x2587)
    gap = Inches(0.3)
    total_w = img_w * 2 + gap
    left0 = CONT_LEFT + (CONT_W - total_w) / 2
    top0 = Inches(1.55)

    add_image(slide, FIG_PLAND_K3_WAD, left0, top0, height=img_h)
    add_image(slide, FIG_PLAND_K3_D0, left0 + img_w + gap, top0, height=img_h)

    add_textbox(slide, "Trainable delays (wad_mlp) — 92.7% acc",
                left0, top0 + img_h + Inches(0.05), img_w, Inches(0.4),
                font_size=13, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "No delays (d0_mlp) — ~77% acc",
                left0 + img_w + gap, top0 + img_h + Inches(0.05), img_w, Inches(0.4),
                font_size=13, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    return slide


def slide9_diagnostic(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Diagnostic Panel Example — Topology 3, K=3, wad_mlp")

    img_w = Inches(7.2)
    img_h = Inches(5.84)  # aspect-preserved (1765x1433)
    left = CONT_LEFT + (CONT_W - img_w) / 2
    add_image(slide, FIG_DIAGNOSTIC_T3, left, Inches(1.3), width=img_w)

    add_textbox(slide,
                "One run, one figure: spike raster (sorted by first-spike time, coloured by sub-window), "
                "delay heatmap, and readout trace — this is what \"routing\" looks like inside the network.",
                CONT_LEFT, Inches(6.95), CONT_W, Inches(0.5),
                font_size=12, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)
    return slide


def slide10_energy(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Energy / Computation Trade-off")

    img_w = Inches(7.6)
    img_h = Inches(5.37)  # aspect-preserved (1078x761)
    left = CONT_LEFT + (CONT_W - img_w) / 2
    add_image(slide, FIG_ENERGY_CHART, left, Inches(1.15), width=img_w)

    add_textbox(slide,
                "Same energy budget (spikes/trial) buys ~3-4× more computations with async delays vs. "
                "synchronous (K independent trials) baseline — up to -73% energy for the same K.",
                CONT_LEFT, Inches(6.6), CONT_W, Inches(0.7),
                font_size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    return slide


def slide11_conclusions(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Conclusions: Delays Are Necessary in Every Shared-Channel Topology")

    claims = [
        ("1.  Structural necessity tracks input-side collision, not output multiplicity",
         "T1 (output multiplicity, no input collision): delays optional, pure alignment (+14-19pp flat).\n"
         "Plan D / T2 / T3 (input collision present): delays structurally necessary, d=0 collapses to label prior."),
        ("2.  Best readout type flips with topology — there's no universal answer",
         "K-parallel-outputs (T1): Linear wins — heads already separate, MLP adds harmful coupling.\n"
         "K-to-1 aggregation (T2/T3) or K-sequential-routing (Plan D): MLP wins — needs nonlinear counting/decoding."),
        ("3.  Independent bottlenecks compound multiplicatively under composition",
         "Routing (Plan D, K@90%=3) + mixed-op competition (Dir C, K@90%=3) + aggregation (T2, K=3) "
         "individually each reach K=3 — but T3 (all three combined) only reaches K=1."),
    ]
    colors = [
        (RGBColor(0xE8, 0xF0, 0xFE), C_IC),
        (RGBColor(0xFF, 0xF3, 0xE0), C_ORANGE),
        (RGBColor(0xD6, 0xE8, 0xD6), C_GREEN),
    ]
    top_start = CONT_TOP
    box_h = Inches(1.5)
    gap = Inches(0.12)
    for idx, (title, body) in enumerate(claims):
        bg, border = colors[idx]
        add_callout_box(slide, f"▶  {title}\n    {body}",
                        CONT_LEFT, top_start + idx * (box_h + gap),
                        CONT_W, box_h, bg_color=bg, border_color=border, font_size=13)

    add_textbox(slide,
                "Plus: depth (L2-h50h50) extends single-op Max K@90% to 4 — the project's best result so far, "
                "with the budget-confound from Sections 13-14 now resolved.",
                CONT_LEFT, Inches(6.1), CONT_W, Inches(0.6),
                font_size=14, bold=True, color=C_NAVY)
    return slide


def slide12_next(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Next Steps")

    bullets = [
        ("Locate the true T1 ceiling: sweep K_ops ∈ {5,6,8} using all 8 canonical boolean ops "
         "(current K_ops≤4 shows no ceiling at h=50).", 0),
        ("Run the missing L1-h100 single-op NAND control to fully disentangle Section 27's "
         "depth effect from raw neuron count.", 0),
        ("Write up the Pattern Taxonomy as a dedicated section/table for the paper — "
         "the 4-row comparison (Plan D / T1 / T2 / T3) is a clean, self-contained contribution.", 0),
        ("Investigate whether T3's compounding bottleneck is fixable with the width/data levers "
         "that worked for Step 3 Direction A/C (more data, larger h), or is a hard ceiling.", 0),
        ("Decide whether a 3rd aggregation function (e.g. parity / XOR-of-results) is worth adding "
         "to T2 for a fuller picture before finalising the taxonomy.", 0),
    ]
    add_bullet_box(slide, bullets, CONT_LEFT, CONT_TOP, CONT_W, Inches(3.8), font_size=15)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    print("Building Step 4 update slides...")
    slide1_title(prs); print("  [1/12] Title")
    slide2_recap(prs); print("  [2/12] Recap")
    slide3_taxonomy(prs); print("  [3/12] Taxonomy")
    slide4_t1(prs); print("  [4/12] Topology 1")
    slide5_t2(prs); print("  [5/12] Topology 2")
    slide6_t3(prs); print("  [6/12] Topology 3")
    slide7_section27(prs); print("  [7/12] Section 27 Depth Ablation")
    slide8_mechanism(prs); print("  [8/12] Mechanism (spike-flow)")
    slide9_diagnostic(prs); print("  [9/12] Diagnostic Panel")
    slide10_energy(prs); print(" [10/12] Energy Chart")
    slide11_conclusions(prs); print(" [11/12] Conclusions")
    slide12_next(prs); print(" [12/12] Next Steps")

    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}")


if __name__ == "__main__":
    main()
