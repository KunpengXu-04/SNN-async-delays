"""
make_presentation_reset.py
Supervisor progress report covering the 2026-07-11 protocol reset: which claims
from the Steps 1-4 deck are retired, why, and the replacement programme.

Run: python make_presentation_reset.py
Requires: pip install python-pptx

Evidence sources (do not edit numbers here without re-checking these):
  - docs/CLAIMS_LEDGER.md          -- claim status labels
  - docs/EXPERIMENT_LOG_V2.md      -- protocol reset onward (2026-07-11+)
  - docs/PUBLICATION_ROADMAP.md    -- positioning, decision tree, open decisions
NOTE: docs/EXPERIMENT_LOG.md (Sections 1-35) is the *historical* narrative and
      contains superseded claims; it is not an evidence source for this deck.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
PRESENTATION_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# Never overwrite the manually maintained current deck (see presentation/README.md).
OUT_PATH = os.path.join(PRESENTATION_ROOT, "generated", "presentation_protocol_reset.pptx")

W = Inches(13.33)
H = Inches(7.5)

# Palette: inherits the existing deck's Imperial navy, adds status colours.
C_NAVY  = RGBColor(0x1F, 0x38, 0x64)
C_IC    = RGBColor(0x00, 0x3E, 0x74)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
C_ALT   = RGBColor(0xEE, 0xF2, 0xF7)
C_RED   = RGBColor(0xB3, 0x26, 0x1E)   # retired / failed
C_TEAL  = RGBColor(0x0F, 0x76, 0x6E)   # passed / established
C_AMBER = RGBColor(0xB4, 0x53, 0x09)   # pending / exploratory
C_GREY  = RGBColor(0x5B, 0x64, 0x70)
C_FAINT = RGBColor(0xF6, 0xF8, 0xFA)

FONT = "Calibri"

TITLE_TOP, TITLE_LEFT, TITLE_W, TITLE_H = Inches(0.22), Inches(0.45), Inches(12.4), Inches(0.62)
CONT_LEFT, CONT_W = Inches(0.45), Inches(12.43)
CONT_TOP = Inches(1.30)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, fill, line=None, line_w=1.0,
             shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_text(slide, text, left, top, width, height, size=16, bold=False,
             italic=False, color=C_BLACK, align=PP_ALIGN.LEFT, anchor=None,
             space_after=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def add_title_bar(slide, title, kicker=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.08), C_NAVY)
    add_text(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H,
             size=27, bold=True, color=C_WHITE)
    if kicker:
        add_text(slide, kicker, TITLE_LEFT, Inches(0.76), TITLE_W, Inches(0.26),
                 size=12, color=RGBColor(0xB8, 0xC7, 0xDE))


def add_chip(slide, text, left, top, color, width=Inches(1.05), height=Inches(0.26),
             size=9.5):
    """Status pill -- the deck's repeated motif."""
    sh = add_rect(slide, left, top, width, height, color,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    sh.adjustments[0] = 0.5
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = FONT
    return sh


def add_bullets(slide, items, left, top, width, height, size=14, color=C_BLACK,
                space=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def make_table(slide, headers, rows, left, top, width, height, col_widths=None,
               size=12, header_size=12, row_colors=None):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers),
                                 left, top, width, height).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)

    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = C_NAVY
        c.margin_left = c.margin_right = Inches(0.06)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(header_size)
        r.font.color.rgb = C_WHITE
        r.font.name = FONT

    for i, row in enumerate(rows):
        bg = row_colors[i] if row_colors else (C_ALT if i % 2 else C_WHITE)
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            c.margin_left = c.margin_right = Inches(0.06)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(size)
            r.font.color.rgb = C_BLACK
            r.font.name = FONT
    return tbl


def add_callout(slide, text, left, top, width, height, accent=C_IC,
                bg=C_FAINT, size=14, bold=False):
    add_rect(slide, left, top, width, height, bg, line=accent, line_w=1.25)
    add_text(slide, text, left + Inches(0.16), top + Inches(0.11),
             width - Inches(0.32), height - Inches(0.22),
             size=size, bold=bold, color=C_BLACK)


def add_card(slide, title, body, left, top, width, height, accent, chip=None):
    """Card with an optional status pill stacked ABOVE the title.

    The pill must not share a line with the title -- long titles collide with a
    right-aligned pill and get visually truncated.
    """
    add_rect(slide, left, top, width, height, C_WHITE, line=RGBColor(0xD6, 0xDD, 0xE5))
    y = top + Inches(0.16)
    if chip:
        add_chip(slide, chip, left + Inches(0.16), y, accent, width=Inches(1.04))
        y = top + Inches(0.56)
    add_text(slide, title, left + Inches(0.16), y, width - Inches(0.32), Inches(0.5),
             size=14, bold=True, color=accent)
    body_top = y + Inches(0.44)
    add_bullets(slide, body, left + Inches(0.16), body_top,
                width - Inches(0.32), height - (body_top - top) - Inches(0.16),
                size=11.5, color=C_GREY, space=5)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = new_slide(prs)
    add_rect(s, Inches(0), Inches(0), W, H, C_NAVY)
    add_text(s, "Protocol Reset", Inches(0.9), Inches(2.25), Inches(11.5), Inches(0.85),
             size=48, bold=True, color=C_WHITE)
    add_text(s, "Correcting the delay-advantage claim, and the programme that replaces it",
             Inches(0.9), Inches(3.18), Inches(11.5), Inches(0.5),
             size=21, color=RGBColor(0xB8, 0xC7, 0xDE))
    add_rect(s, Inches(0.9), Inches(3.95), Inches(1.5), Inches(0.045), C_TEAL)
    add_bullets(s, [
        "Progress report covering 2026-06-29 (last deck) to 2026-07-16",
        "Supersedes the Max K / energy claims in the Steps 1-4 deck",
    ], Inches(0.9), Inches(4.3), Inches(11.0), Inches(0.9), size=15,
        color=RGBColor(0xD5, 0xDF, 0xEE), space=4)
    add_text(s, "SNN Async Delays Project   |   Imperial College London",
             Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.3),
             size=13, color=RGBColor(0x8F, 0xA5, 0xC4))
    s.notes_slide.notes_text_frame.text = (
        "Frame this up front: this is a correction, not an incremental update. "
        "The headline numbers in the last deck do not survive re-evaluation. "
        "Say plainly that I found the problem myself, through a protocol audit, "
        "and that the audit is the contribution."
    )
    return s


def slide_summary(prs):
    s = new_slide(prs)
    add_title_bar(s, "Executive summary",
                  "Three numbers that changed between the last deck and today")

    stats = [
        ("K = 2", "Max K that survives\nre-evaluation",
         "Last deck claimed 3\n(later 4-5). 177 checkpoints\nre-scored; K>=3 fails\nthrough h=150.", C_RED),
        ("1.000 = 1.000", "d0 vs WAD under a\nneutral readout",
         "Under late-window MLP it\nread 0.490 vs 0.727. The gap\nwas the measurement\ninterface, not the network.", C_RED),
        ("1 / 5", "Seeds where WAD beat\nan optimized scalar delay",
         "Preregistered gate needed\n>=4/5 and mean >= +.03.\nActual mean: -.021.\nGate failed decisively.", C_RED),
    ]
    x = CONT_LEFT
    cw = Inches(3.95)
    gap = Inches(0.29)
    for big, label, note, col in stats:
        add_rect(s, x, CONT_TOP, cw, Inches(2.75), C_WHITE,
                 line=RGBColor(0xD6, 0xDD, 0xE5))
        add_text(s, big, x + Inches(0.2), CONT_TOP + Inches(0.22), cw - Inches(0.4),
                 Inches(0.62), size=34, bold=True, color=col)
        add_text(s, label, x + Inches(0.2), CONT_TOP + Inches(0.92), cw - Inches(0.4),
                 Inches(0.55), size=13, bold=True, color=C_BLACK)
        add_text(s, note, x + Inches(0.2), CONT_TOP + Inches(1.55), cw - Inches(0.4),
                 Inches(1.05), size=11, color=C_GREY)
        x += cw + gap

    add_callout(s,
                "What happened:  a July protocol audit found that the readout interface used in every prior "
                "experiment censors the no-delay baseline's activity. Once the interface is made neutral, the "
                "delay advantage disappears -- and an optimized single scalar delay outperforms learned "
                "heterogeneous delays. The positive learned-delay programme has been rejected on its own "
                "preregistered gate.",
                CONT_LEFT, Inches(4.35), CONT_W, Inches(1.15), accent=C_RED, size=14)

    add_text(s, "The work since the last deck is therefore a correction plus a rebuilt, "
                "preregistered mechanism programme -- not a set of new performance results.",
             CONT_LEFT, Inches(5.72), CONT_W, Inches(0.4), size=14, bold=True, color=C_IC)

    add_text(s, "Sources: docs/CLAIMS_LEDGER.md, docs/EXPERIMENT_LOG_V2.md (2026-07-11 onward)",
             CONT_LEFT, Inches(6.95), CONT_W, Inches(0.25), size=10, color=C_GREY)
    s.notes_slide.notes_text_frame.text = (
        "Do not rush this slide. These three numbers are the whole talk. "
        "If the supervisor only remembers one thing, it should be the middle one: "
        "d0 and WAD are both perfect under a fair readout."
    )
    return s


def slide_confound(prs):
    s = new_slide(prs)
    add_title_bar(s, "Root cause: the readout window manufactured the result",
                  "XOR readout-interface pilot v1  |  36 preregistered cells  |  K=2, h=50, seeds {0,1,42}")

    add_text(s, "Worst-query accuracy, mean +/- SD over 3 seeds",
             CONT_LEFT, CONT_TOP, Inches(7.6), Inches(0.28), size=12, bold=True, color=C_GREY)

    rows = [
        ["late_window  +  MLP", "0.490 +/- .017", "0.727 +/- .132", "+0.237"],
        ["all_time  +  MLP", "1.000 +/- .000", "1.000 +/- .000", "0.000"],
        ["all_time  +  linear", "0.850 +/- .090", "0.877 +/- .054", "+0.027"],
        ["time_binned (cost-inflated)", "1.000", "1.000", "0.000"],
    ]
    make_table(s, ["Observation interface", "d0", "WAD", "Gap"], rows,
               CONT_LEFT, Inches(1.68), Inches(7.6), Inches(1.85),
               col_widths=[3.0, 1.5, 1.5, 1.0], size=12,
               row_colors=[RGBColor(0xFA, 0xE8, 0xE6), C_WHITE, C_ALT, C_WHITE])

    add_callout(s,
                "The historical late-window readout only observes the final window. It discards the early "
                "activity that d0 relies on -- so d0 is scored at chance by construction, and WAD's "
                "\"necessity\" is an artifact of the measurement, not a property of delays.",
                CONT_LEFT, Inches(3.72), Inches(7.6), Inches(1.18), accent=C_RED, size=13.5)

    add_text(s, "Consequences now locked into the protocol", CONT_LEFT, Inches(5.08),
             Inches(7.6), Inches(0.3), size=13, bold=True, color=C_IC)
    add_bullets(s, [
        "all_time is the primary, neutral observation interface.",
        "late_window is demoted to an alignment / censoring diagnostic.",
        "time_binned is an upper bound only -- it reaches 1.000 but raises decoder",
        "     storage 2,852 -> 23,152 and decoder MACs 2,600 -> 22,800 (~9x).",
        "Plan D is renamed a sequential delayed-retention / routing task -- not",
        "     \"proven temporal multiplexing\".",
    ], CONT_LEFT, Inches(5.44), Inches(7.6), Inches(1.5), size=12.5, color=C_BLACK, space=3)

    # Right column: the interpretation panel
    px = Inches(8.35)
    pw = Inches(4.53)
    add_rect(s, px, CONT_TOP, pw, Inches(4.35), C_NAVY)
    add_text(s, "Why this is the contribution", px + Inches(0.24), CONT_TOP + Inches(0.24),
             pw - Inches(0.48), Inches(0.35), size=16, bold=True, color=C_WHITE)
    add_bullets(s, [
        "This confound is not unique to us. Any delay-SNN paper that reads out "
        "only at the end of a trial can report the same illusory advantage.",
        "",
        "We found it by building the neutral interface and re-running the "
        "comparison -- not by inspecting rasters.",
        "",
        "That makes \"observation-induced delay advantage\" a publishable "
        "methodological result in its own right, independent of whether "
        "delays ultimately win.",
    ], px + Inches(0.24), CONT_TOP + Inches(0.78), pw - Inches(0.48), Inches(3.3),
        size=12.5, color=RGBColor(0xD5, 0xDF, 0xEE), space=4)

    add_chip(s, "ROOT CAUSE", px + Inches(0.24), Inches(5.82), C_RED, width=Inches(1.25))
    add_text(s, "Everything on the next slide follows from this one finding.",
             px + Inches(1.62), Inches(5.83), Inches(2.9), Inches(0.3),
             size=11.5, italic=True, color=C_GREY)

    s.notes_slide.notes_text_frame.text = (
        "This is the most important slide. Walk the table row by row. "
        "Row 1 is what we used to do; row 2 is the same models under a fair readout. "
        "Expect the question 'was the old readout wrong or just different?' -- answer: "
        "it is a legitimate interface, but it cannot support a capacity claim, only an "
        "alignment claim. That distinction is the whole point."
    )
    return s


def slide_retired(prs):
    s = new_slide(prs)
    add_title_bar(s, "Claims retired from the Steps 1-4 deck",
                  "Status labels are taken verbatim from docs/CLAIMS_LEDGER.md")

    rows = [
        ["Max K@90% = 3  (later 4 / 5)",
         "Full re-evaluation of all 177 NAND-burst checkpoints: requiring every seed to clear "
         "worst-query >= 90%, WAD passes only K=1 and K=2. No h through 150 passes K>=3."],
        ["Delays are the sole necessary mechanism\nfor shared-channel temporal routing",
         "Interface artifact. Under all_time MLP both d0 and WAD reach 1.000. The d0 failure was "
         "produced by late-window censoring."],
        ["19% fewer spikes  /  73% energy saving",
         "Historical Step-1 numbers disagree with canonical d0 results, and hidden-spike count is "
         "not an energy model. No scalar energy claim without declared hardware."],
        ["Delays yield a fixed-resource\ncapacity gain",
         "T, output width, MLP size and delay range all grow with K. The resources were never "
         "matched, so the comparison cannot support the claim."],
        ["The network multiplexes\nrather than aligns",
         "Never tested. The late-window design is fully compatible with alignment alone; no causal "
         "intervention distinguished them."],
        ["Burst reveals a timing mechanism;\nhomeostasis breaks the sparsity tradeoff",
         "Downgraded, not deleted. Both rest on selected, low-seed observations and need a "
         "predeclared multi-seed study."],
    ]
    statuses = [
        ("RETIRED", C_RED), ("RETIRED", C_RED), ("RETIRED", C_RED),
        ("RETIRED", C_RED), ("NOT TESTED", C_AMBER), ("EXPLORATORY", C_AMBER),
    ]

    y = Inches(1.32)
    rh = Inches(0.88)
    for (claim, why), (label, col) in zip(rows, statuses):
        add_rect(s, CONT_LEFT, y, CONT_W, rh - Inches(0.06), C_WHITE,
                 line=RGBColor(0xE0, 0xE5, 0xEB))
        add_text(s, claim, CONT_LEFT + Inches(0.16), y + Inches(0.12), Inches(3.35),
                 Inches(0.62), size=12, bold=True, color=C_BLACK)
        add_chip(s, label, CONT_LEFT + Inches(3.66), y + Inches(0.25), col, width=Inches(1.12))
        add_text(s, why, CONT_LEFT + Inches(4.95), y + Inches(0.13), Inches(7.3),
                 Inches(0.62), size=11.5, color=C_GREY)
        y += rh

    add_text(s, "Nothing here was retracted by an outside reviewer -- every line was found and "
                "recorded by the project's own audit and ledger.",
             CONT_LEFT, Inches(6.72), CONT_W, Inches(0.32), size=12.5, italic=True, color=C_IC)
    s.notes_slide.notes_text_frame.text = (
        "Don't linger defensively. State each retirement flatly and move on. "
        "The bottom line is the message: the ledger caught these, which is exactly "
        "what a claims ledger is for."
    )
    return s


def slide_gate1(prs):
    s = new_slide(prs)
    add_title_bar(s, "Gate 1 failed: an optimized scalar delay beats learned delays",
                  "XOR delay control matrix v1  |  60 cells  |  K=3, N=35, T=40, all-time linear  |  5 paired seeds")

    add_text(s, "Primary endpoint -- worst-query accuracy", CONT_LEFT, CONT_TOP,
             Inches(7.1), Inches(0.28), size=12.5, bold=True, color=C_GREY)
    make_table(s, ["Condition", "Worst-query", "Rank"], [
        ["Optimized shared scalar delay", ".699 +/- .019", "1 (best)"],
        ["WAD  (learned heterogeneous)", ".678 +/- .036", "2"],
    ], CONT_LEFT, Inches(1.66), Inches(7.1), Inches(0.92),
        col_widths=[3.6, 1.9, 1.6], size=12.5,
        row_colors=[RGBColor(0xE3, 0xF2, 0xF0), RGBColor(0xFA, 0xE8, 0xE6)])

    add_text(s, "Paired WAD - scalar, per seed:      -.028    -.060    -.030    -.036    +.050",
             CONT_LEFT, Inches(2.74), Inches(7.1), Inches(0.3), size=13, bold=True, color=C_RED)
    add_text(s, "mean  -.0208     |     positive in 1 / 5 seeds",
             CONT_LEFT, Inches(3.04), Inches(7.1), Inches(0.3), size=12.5, color=C_GREY)

    add_callout(s,
                "Preregistered gate:  mean paired advantage >= +.03, positive in >= 4/5 seeds, no resource "
                "domination.\nResult:  mean -.021, positive in 1/5.  Decision recorded: reject the positive "
                "learned-delay superiority programme.",
                CONT_LEFT, Inches(3.5), Inches(7.1), Inches(1.02), accent=C_RED, size=12.5)

    add_text(s, "The result is not an optimization failure", CONT_LEFT, Inches(4.72),
             Inches(7.1), Inches(0.3), size=13, bold=True, color=C_IC)
    add_bullets(s, [
        "A 60-cell optimization audit (Stage A + B) tried d_max, delay LR, noisy init,",
        "     warm-up and alternating schedules. Mean WAD worst-query:",
        "     baseline .676  >  noisy-init .654  =  low-LR .654  >  alternating .641",
        "     >  d10 .638  >  warm-up .624  >  high-LR .607.",
        "Every intervention scored BELOW baseline. Delay gradients were finite and",
        "     nonzero throughout; saturation was zero; delays did move (~.99 steps).",
    ], CONT_LEFT, Inches(5.06), Inches(7.1), Inches(1.5), size=12, color=C_BLACK, space=3)

    # Right panel
    px, pw = Inches(7.85), Inches(5.03)
    add_rect(s, px, CONT_TOP, pw, Inches(2.5), C_WHITE, line=RGBColor(0xD6, 0xDD, 0xE5))
    add_text(s, "Secondary endpoints agree", px + Inches(0.18), CONT_TOP + Inches(0.15),
             pw - Inches(0.36), Inches(0.3), size=13, bold=True, color=C_IC)
    add_bullets(s, [
        "MLP readout:  fixed heterogeneous .809  >  scalar .780",
        "     >  d0 .775  >  WAD .770      (WAD ranks last)",
        "Stress point (K=4, N=50):  scalar .652  >  WAD .644",
        "Delay shuffle degrades WAD in 15/15 cells (mean -.116)",
        "     -- evidence of co-adapted placement, NOT of superiority.",
    ], px + Inches(0.18), CONT_TOP + Inches(0.55), pw - Inches(0.36), Inches(1.8),
        size=11.5, color=C_GREY, space=4)

    add_rect(s, px, Inches(4.0), pw, Inches(2.55), C_FAINT, line=C_AMBER, line_w=1.25)
    add_chip(s, "STILL OPEN", px + Inches(0.18), Inches(4.16), C_AMBER, width=Inches(1.15))
    add_text(s, "The one control that is still missing",
             px + Inches(0.18), Inches(4.52), pw - Inches(0.36), Inches(0.3),
             size=13, bold=True, color=C_BLACK)
    add_bullets(s, [
        "The fixed-heterogeneous control is uniform 0-30 (mean ~15), while WAD "
        "concentrates at 1-9 (mean ~3.5). It is NOT a distribution-matched control.",
        "",
        "So \"delay structure vs delay learning\" is genuinely unresolved -- a matched "
        "bank or a cross-fitted WAD multiset is required before that question can "
        "be closed either way.",
    ], px + Inches(0.18), Inches(4.9), pw - Inches(0.36), Inches(1.5),
        size=11.5, color=C_GREY, space=5)

    s.notes_slide.notes_text_frame.text = (
        "Key line: beating d0 was never the bar; the bar is beating the strongest "
        "non-learned control, and a single trainable scalar delay is that control. "
        "Flag the still-open box honestly -- the supervisor will likely ask whether "
        "the fixed control was fair, and the answer is: not yet, and that is logged."
    )
    return s


def slide_more_negatives(prs):
    s = new_slide(prs)
    add_title_bar(s, "Three further independent negatives",
                  "Different tasks, different interfaces, same conclusion -- this is not one bad experiment")

    cards = [
        ("Counterbalanced temporal routing", C_RED, [
            "45 cells, seeds {7,19,73}, every operation in every temporal position.",
            "",
            "Worst operation x position accuracy is exactly .500 for d0, scalar, "
            "fixed-matched AND WAD. Fixed-full: .500/.486/.482.",
            "",
            "No condition passes the .55 per-seed / .60 mean routing floor.",
            "",
            "Diagnosis: position 2 is the barrier, not operation identity. WAD learns "
            "NOR early but fails NOR and XOR late.",
        ]),
        ("Spatial vs temporal Pareto scaffold", C_RED, [
            "160 cells, K=2 XOR, MLP scaffold with output-spike conversion removed.",
            "",
            "WAD worst-balanced = .500 in ALL 32 of its cells; window-1 hidden spikes "
            "are exactly ZERO. Its delays never separate by query.",
            "",
            "A fixed oracle schedule DOES work (24 vs 48 hidden at T=18) -- so the "
            "interface can express routing; WAD just fails to learn it.",
            "",
            "But the oracle saves nothing on dense MACs, events, delay storage or "
            "decoder MACs -- and costs 5x the delay-buffer memory.",
        ]),
        ("Optimization rescue audit", C_RED, [
            "60 cells across Stage A (threshold viability) and Stage B (7 variants).",
            "",
            "Threshold 0.3 selected by a predeclared activity/gradient rule "
            "(mean activity 10.07 vs target 10) -- explicitly NOT by accuracy.",
            "",
            "No variant passed the rescue gate. All 6 interventions scored below the "
            ".676 baseline.",
            "",
            "Rules out the obvious objection: WAD is not simply under-trained or "
            "badly tuned.",
        ]),
    ]

    x = CONT_LEFT
    cw = Inches(3.95)
    gap = Inches(0.29)
    for title, accent, body in cards:
        add_card(s, title, body, x, CONT_TOP, cw, Inches(4.55), accent, chip="NEGATIVE")
        x += cw + gap

    add_callout(s,
                "Taken together: the delay advantage does not survive a neutral readout, a stronger control, "
                "a counterbalanced task, a different scaffold, or a tuning rescue. The convergence of five "
                "independent lines is itself the evidence.",
                CONT_LEFT, Inches(6.05), CONT_W, Inches(0.9), accent=C_IC, size=13.5)
    s.notes_slide.notes_text_frame.text = (
        "Purpose of this slide is to pre-empt 'maybe you just got unlucky once'. "
        "Five independent lines, all negative. Note the Pareto card's silver lining: "
        "a fixed oracle schedule works, which tells us the target is learnable in "
        "principle -- our method just doesn't find it."
    )
    return s


def slide_ladder(prs):
    s = new_slide(prs)
    add_title_bar(s, "The replacement programme: a bottom-up mechanism ladder",
                  "Stop asserting performance top-down. Prove one delay can be learned at all -- one mechanism at a time.")

    add_callout(s,
                "Rationale:  every failed experiment above assumed the delay-learning machinery works and "
                "asked whether it helps. The ladder inverts that -- it strips away XOR, labels, weights and "
                "multiple queries until only ONE synapse and ONE event remain, then adds complexity back "
                "only through a passed gate.",
                CONT_LEFT, CONT_TOP, CONT_W, Inches(0.92), accent=C_IC, size=13)

    rows = [
        ["0A", "Can the sigmoid delay parameter recover a declared target at all?",
         "75", "FAILED at the production recipe",
         "'.001'/200 steps moves the delay 0.954 -> 1.139 against target 5 (error 3.86). "
         "'.05' recovers 15/15. So: no hard implementation bug, but the update budget is "
         "inadequate even under DIRECT supervision."],
        ["0B", "Does timing credit survive the real buffer, then one LIF neuron?",
         "180", "FAILED strict gate",
         "Buffer + arrival centroid: 15/15 -- the circular buffer is NOT the obstruction. "
         "But hard-LIF + centroid: 3/15; a hard spike's centroid is piecewise constant, so "
         "the gradient is exactly zero in 8/13 misaligned pairs."],
        ["0C", "Can a soft/global trace give correct bidirectional credit?",
         "360", "PASSED",
         "Selected: sigmoid + soft centroid + Adam '.05' recovers 30/30 across current and "
         "membrane paths, with 13/13 correct initial directions and zero hard spikes."],
        ["0D", "Can that credit place a HARD spike at a declared time?",
         "135", "PASSED",
         "Selected: hard loss + current centroid at lambda=.1 -> 15/15 with 13/13 correct "
         "directions. Hard loss alone is only 10/15 with 5 WRONG directions. Pre-reset "
         "voltage hypothesis rejected (3/15)."],
    ]

    # Layout: [accent | level+cells+pill | question | detail]. The pill lives in the
    # left column so the detail column keeps a full, uncrowded width.
    y = Inches(2.42)
    rh = Inches(1.08)
    for level, q, cells, verdict, detail in rows:
        passed = verdict.startswith("PASSED")
        col = C_TEAL if passed else C_RED
        add_rect(s, CONT_LEFT, y, CONT_W, rh - Inches(0.08), C_WHITE,
                 line=RGBColor(0xE0, 0xE5, 0xEB))
        add_rect(s, CONT_LEFT, y, Inches(0.09), rh - Inches(0.08), col)
        add_text(s, "Level " + level, CONT_LEFT + Inches(0.24), y + Inches(0.11),
                 Inches(1.0), Inches(0.28), size=13, bold=True, color=col)
        add_text(s, cells + " cells", CONT_LEFT + Inches(0.24), y + Inches(0.40),
                 Inches(1.0), Inches(0.24), size=10, color=C_GREY)
        add_chip(s, "PASSED" if passed else "FAILED",
                 CONT_LEFT + Inches(0.24), y + Inches(0.66), col, width=Inches(0.95))
        add_text(s, q, CONT_LEFT + Inches(1.5), y + Inches(0.16), Inches(3.7),
                 Inches(0.7), size=11.5, bold=True, color=C_BLACK)
        add_text(s, detail, CONT_LEFT + Inches(5.42), y + Inches(0.13), Inches(6.85),
                 Inches(0.8), size=11, color=C_GREY)
        y += rh

    add_text(s, "750 deterministic unit cells, every one preregistered with its gate written before the run.",
             CONT_LEFT, Inches(6.92), CONT_W, Inches(0.3), size=12, italic=True, color=C_IC)
    s.notes_slide.notes_text_frame.text = (
        "The story: 0A and 0B are negative but diagnostic -- they tell us exactly which "
        "component was broken (the objective, not the buffer, not the parameterisation). "
        "0C and 0D then build a credit path that provably works. This is the first time "
        "the project has a verified delay-learning mechanism at any scale."
    )
    return s


def slide_level1a(prs):
    s = new_slide(prs)
    add_title_bar(s, "Level 1A: the first real task -- passed, on a deliberately narrow claim",
                  "K=1 XOR bridge  |  Stage I 90 cells + Stage II 85 cells  |  4 -> 16 -> 2 hard-spiking net")

    add_text(s, "What passed", CONT_LEFT, CONT_TOP, Inches(6.1), Inches(0.3),
             size=14, bold=True, color=C_TEAL)
    add_bullets(s, [
        "Stage I (90 cells): the fixed-schedule hard-spike XOR interface passes "
        "10/10 across both schedules and 5 seeds. Selected eta=0, lr_w=.01.",
        "",
        "Stage II (85 cells): with an arrival-centroid scaffold, the learned delay "
        "lands on target in 10/10 cells. Selected lambda=.01, lr_d=.01.",
        "",
        "Final delay range [3.997372, 4.001881] -- maximum error .002628 step, "
        "against a target of 4. Exact hard-spike XOR truth table retained in all 10.",
    ], CONT_LEFT, Inches(1.68), Inches(6.1), Inches(2.3), size=12.5, color=C_BLACK, space=4)

    add_rect(s, CONT_LEFT, Inches(4.05), Inches(6.1), Inches(1.0), C_WHITE, line=C_TEAL, line_w=1.25)
    add_chip(s, "PASSED", CONT_LEFT + Inches(0.16), Inches(4.2), C_TEAL, width=Inches(0.95))
    add_text(s, "First verified joint task + delay learning anywhere in the project.",
             CONT_LEFT + Inches(1.28), Inches(4.22), Inches(4.7), Inches(0.6),
             size=12.5, bold=True, color=C_BLACK)

    # Right: the boundary
    px, pw = Inches(6.85), Inches(6.03)
    add_rect(s, px, CONT_TOP, pw, Inches(3.9), RGBColor(0xFA, 0xE8, 0xE6), line=C_RED, line_w=1.5)
    add_chip(s, "CLAIM BOUNDARY", px + Inches(0.2), CONT_TOP + Inches(0.18), C_RED, width=Inches(1.55))
    add_text(s, "What this is NOT", px + Inches(0.2), CONT_TOP + Inches(0.56),
             pw - Inches(0.4), Inches(0.35), size=16, bold=True, color=C_RED)
    add_bullets(s, [
        "The arrival loss EXPLICITLY ENCODES the oracle delay-4 schedule. "
        "The network is told the answer and checked for whether it can hold it.",
        "",
        "Task-only delay learning (lambda=0) passes 0/10 at lr_d=.01, and 1/10 "
        "at .05. Only 5/10 initial task-gradient directions are even correct.",
        "",
        "Timing-specificity control (d0 forced to target t=15): 0/5, with a "
        "correct-target-time rate of zero in every seed.",
        "",
        "=> This is a scaffold-assisted bridge. It is NOT evidence of "
        "task-derived routing discovery, WAD, multiplexing, generalisation, "
        "compression, or a Pareto law.",
    ], px + Inches(0.2), CONT_TOP + Inches(1.0), pw - Inches(0.4), Inches(2.8),
        size=12, color=C_BLACK, space=4)

    add_callout(s,
                "Read honestly: we can now learn a delay when we tell the network which delay to learn. "
                "Whether the task itself can drive that discovery is exactly what Level 1B onward must test.",
                px, Inches(5.35), pw, Inches(0.85), accent=C_AMBER, size=12.5)

    add_text(s, "Sources: RESULTS_XOR_TASK_BRIDGE_LEVEL1A_STAGE_I.md / _STAGE_II.md",
             CONT_LEFT, Inches(6.95), CONT_W, Inches(0.25), size=10, color=C_GREY)
    s.notes_slide.notes_text_frame.text = (
        "Be scrupulous here. If the supervisor hears 'Level 1A passed' without the "
        "boundary, they will over-read it. The lambda=0 result (0/10) is the honest "
        "headline: without the oracle scaffold, the task cannot teach the delay yet."
    )
    return s


def slide_level1b(prs):
    s = new_slide(prs)
    add_title_bar(s, "Level 1B: preregistered, launch-ready, not yet run",
                  "Does the Level-1A scaffold survive more delay coordinates and a micro-burst input?")

    add_text(s, "Locked design", CONT_LEFT, CONT_TOP, Inches(6.1), Inches(0.3),
             size=14, bold=True, color=C_IC)
    add_bullets(s, [
        "Stage A -- 60 cells:  3 delay granularities (global 1 / per-neuron 16 / "
        "per-synapse 64)  x  2 loss conditions (task-only vs scaffold)  x  2 matched "
        "initial directions  x  5 NEW seeds {607,709,811,919,1021}, disjoint from 1A.",
        "",
        "Stage B -- micro-burst robustness:  events at steps 8 and 9, target still "
        "exactly ONE opponent spike at t=15. Mechanically locked until a fixed-delay-4 "
        "control passes 5/5.",
    ], CONT_LEFT, Inches(1.68), Inches(6.1), Inches(2.0), size=12, color=C_BLACK, space=4)

    add_callout(s,
                "The gate is coordinate-wise, not mean-wise.  A smoke cell reached exact XOR with a mean "
                "delay of 3.99996 and still FAILED -- its worst coordinate was .378 off and coverage was "
                ".859. A mean cannot hide a bad coordinate.",
                CONT_LEFT, Inches(3.82), Inches(6.1), Inches(1.15), accent=C_IC, size=12)

    add_text(s, "Status", CONT_LEFT, Inches(5.15), Inches(6.1), Inches(0.3),
             size=14, bold=True, color=C_IC)
    add_bullets(s, [
        "88/88 project tests pass; dry-run expands to exactly 60 / 10 / 60 unique paths.",
        "6 smoke cells validate code paths only -- explicitly invalid for claims.",
        "0 formal cells launched.",
    ], CONT_LEFT, Inches(5.5), Inches(6.1), Inches(1.0), size=12, color=C_BLACK, space=4)

    # Placeholder panel
    px, pw = Inches(6.85), Inches(6.03)
    ph = Inches(5.1)
    add_rect(s, px, CONT_TOP, pw, ph, C_FAINT, line=C_GREY, line_w=1.5)
    add_chip(s, "PENDING", px + Inches(0.2), CONT_TOP + Inches(0.2), C_AMBER, width=Inches(1.0))
    add_text(s, "RESULTS", px, CONT_TOP + Inches(1.55), pw, Inches(0.6),
             size=40, bold=True, color=RGBColor(0xC3, 0xCB, 0xD4), align=PP_ALIGN.CENTER)
    add_text(s, "to be added once Stage A completes", px, CONT_TOP + Inches(2.25), pw,
             Inches(0.4), size=15, color=C_GREY, align=PP_ALIGN.CENTER)
    add_text(s, "Reserved for:  per-granularity pass counts (x/10)  |  worst-coordinate "
                "delay error  |  coverage  |  task-only vs scaffold contrast",
             px + Inches(0.55), CONT_TOP + Inches(3.1), pw - Inches(1.1), Inches(0.8),
             size=11.5, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

    s.notes_slide.notes_text_frame.text = (
        "Say explicitly that this is unrun and that the gate was written before any "
        "data exists. If the supervisor wants to influence the design, this is the "
        "moment -- after launch it is frozen."
    )
    return s


def slide_positioning(prs):
    s = new_slide(prs)
    add_title_bar(s, "Where the project actually stands",
                  "From docs/PUBLICATION_ROADMAP.md -- the authoritative forward plan")

    add_rect(s, CONT_LEFT, CONT_TOP, CONT_W, Inches(0.95), C_NAVY)
    add_text(s, "\"The defensible project is not currently 'delays enable temporal multiplexing' "
                "or 'delays save energy.'\"",
             CONT_LEFT + Inches(0.25), CONT_TOP + Inches(0.14), CONT_W - Inches(0.5),
             Inches(0.7), size=15, italic=True, bold=True, color=C_WHITE)
    add_text(s, "-- PUBLICATION_ROADMAP.md, Executive decision",
             CONT_LEFT + Inches(0.25), CONT_TOP + Inches(0.62), CONT_W - Inches(0.5),
             Inches(0.25), size=10.5, color=RGBColor(0x8F, 0xA5, 0xC4))

    add_text(s, "The question the evidence can actually support", CONT_LEFT, Inches(2.5),
             Inches(6.1), Inches(0.3), size=13.5, bold=True, color=C_IC)
    add_callout(s,
                "Under a neutral observation interface and an explicitly matched resource vector, do "
                "learnable heterogeneous delays improve the reliability-cost Pareto frontier over d0, "
                "optimized scalar, and fixed heterogeneous controls?",
                CONT_LEFT, Inches(2.85), Inches(6.1), Inches(1.25), accent=C_IC, size=13)

    add_text(s, "Note what changed: \"do delays help?\" became \"do LEARNED delays beat the best "
                "NON-learned delay, at matched cost, under a fair readout?\" Beating d0 is no longer a result.",
             CONT_LEFT, Inches(4.25), Inches(6.1), Inches(0.7), size=12, italic=True, color=C_GREY)

    add_text(s, "Publication decision tree -- which branch we are on", CONT_LEFT, Inches(5.1),
             Inches(6.1), Inches(0.3), size=13.5, bold=True, color=C_IC)
    branches = [
        ("WAD beats all matched controls + survives scaling & an external task", "NOT SUPPORTED", C_RED),
        ("Advantage only under late-window / vanishes at matched decoder info", "WE ARE HERE", C_TEAL),
        ("WAD ties fixed heterogeneous, both beat d0", "UNRESOLVED", C_AMBER),
        ("No robust benefit after controls", "POSSIBLE", C_GREY),
    ]
    y = Inches(5.48)
    for text, label, col in branches:
        here = label == "WE ARE HERE"
        add_rect(s, CONT_LEFT, y, Inches(6.1), Inches(0.38),
                 RGBColor(0xE3, 0xF2, 0xF0) if here else C_WHITE,
                 line=col if here else RGBColor(0xE0, 0xE5, 0xEB),
                 line_w=1.5 if here else 1.0)
        add_text(s, text, CONT_LEFT + Inches(0.12), y + Inches(0.09), Inches(4.45),
                 Inches(0.26), size=10.5, bold=here, color=C_BLACK)
        add_chip(s, label, CONT_LEFT + Inches(4.72), y + Inches(0.06), col,
                 width=Inches(1.25), height=Inches(0.25), size=8.5)
        y += Inches(0.42)

    # Right panel
    px, pw = Inches(6.85), Inches(6.03)
    add_rect(s, px, Inches(2.5), pw, Inches(4.05), C_WHITE, line=RGBColor(0xD6, 0xDD, 0xE5))
    add_text(s, "Venue reality", px + Inches(0.2), Inches(2.66), pw - Inches(0.4),
             Inches(0.3), size=14, bold=True, color=C_IC)
    add_bullets(s, [
        "Current evidence = an exploratory / negative-result study. That is a "
        "workshop paper or a thesis chapter, not a conference method paper.",
        "",
        "A serious conference submission requires Phases 1-5: causal controls, "
        "a K/N/T response surface, mechanism interventions, encoding robustness, "
        "and one external temporal benchmark (SHD or SSC).",
        "",
        "\"We learn delays\" is explicitly NOT novel -- DCLS and the axonal-delay "
        "literature already do it. Novelty has to come from the controlled "
        "evaluation and the demonstrated failure mode.",
        "",
        "The honest asset we now hold: a reusable protocol that exposes "
        "observation-window and decoder-resource confounding in delay-SNN "
        "evaluation. That is a real contribution, and it is publishable.",
    ], px + Inches(0.2), Inches(3.02), pw - Inches(0.4), Inches(3.4), size=11.5,
        color=C_GREY, space=4)

    s.notes_slide.notes_text_frame.text = (
        "This is where I need the supervisor to engage rather than just receive. "
        "The pivotal sentence is the last bullet: the protocol is the asset. "
        "If they want a positive method paper instead, that is Phases 1-5 and "
        "months of compute -- which is the next slide's question."
    )
    return s


def slide_decisions(prs):
    s = new_slide(prs)
    add_rect(s, Inches(0), Inches(0), W, H, C_NAVY)
    add_text(s, "Decisions I need from you", Inches(0.9), Inches(0.62), Inches(11.5),
             Inches(0.6), size=34, bold=True, color=C_WHITE)
    add_text(s, "Recorded as open in PUBLICATION_ROADMAP.md Sec. 16 -- these gate Phase 2, "
                "and I should not pick them unilaterally",
             Inches(0.9), Inches(1.28), Inches(11.5), Inches(0.35), size=14,
             color=RGBColor(0xB8, 0xC7, 0xDE))

    items = [
        ("1", "Method, or methodology?",
         "The evidence now favours an evaluation/protocol contribution over a learned-delay method. "
         "Does the thesis follow the evidence, or do we spend another cycle trying to rescue the "
         "method claim? This is the fork everything else hangs on."),
        ("2", "Is the matched-distribution control worth one more cycle?",
         "It is the single missing control that could still separate \"delay structure\" from "
         "\"delay learning\". Cheap relative to Phases 1-5 -- but it can only narrow a negative, "
         "not create a positive."),
        ("3", "Submission window and venue class?",
         "Drives whether we aim at a workshop/thesis chapter now, or commit to Phases 1-5 and an "
         "external benchmark for a conference."),
        ("4", "Compute budget and confirmatory seed count?",
         "Headline confirmation needs >=10 paired seeds. Current results run at 2-5. This is the "
         "main cost driver."),
        ("5", "Is latency T a constrained resource, or just an axis of study?",
         "Changes the entire cost model -- and whether temporal multiplexing can ever be "
         "\"cheaper\" in a meaningful sense."),
    ]
    # Spacing is tuned so row 5 clears the footer -- keep >=0.3" of bottom margin.
    y = Inches(1.88)
    for num, q, detail in items:
        add_rect(s, Inches(0.9), y, Inches(11.53), Inches(0.9),
                 RGBColor(0x2B, 0x47, 0x76))
        add_text(s, num, Inches(1.08), y + Inches(0.2), Inches(0.45), Inches(0.5),
                 size=22, bold=True, color=C_TEAL)
        add_text(s, q, Inches(1.62), y + Inches(0.1), Inches(4.2), Inches(0.7),
                 size=13.5, bold=True, color=C_WHITE)
        add_text(s, detail, Inches(5.95), y + Inches(0.1), Inches(6.35), Inches(0.7),
                 size=11, color=RGBColor(0xB8, 0xC7, 0xDE))
        y += Inches(0.97)

    add_text(s, "Everything above is validation-only. The test split has never been opened.",
             Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.3), size=12,
             italic=True, color=C_TEAL)

    s.notes_slide.notes_text_frame.text = (
        "Close on question 1 and stop talking. It is genuinely their call and it "
        "determines the next three months. Do not fill the silence by proposing "
        "more experiments."
    )
    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_summary(prs)
    slide_confound(prs)
    slide_retired(prs)
    slide_gate1(prs)
    slide_more_negatives(prs)
    slide_ladder(prs)
    slide_level1a(prs)
    slide_level1b(prs)
    slide_positioning(prs)
    slide_decisions(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
