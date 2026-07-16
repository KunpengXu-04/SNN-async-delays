"""
make_presentation_step3.py
Build the missing Step 3 (mixed-op temporal multiplexing) slides, in the same
visual language as make_presentation.py / make_presentation_step4.py — charts
over tables. Output is a standalone deck meant to be merged into
presentation.pptx between the Steps 1-2 conclusions slide and the Step 4 slides.
Run: python make_presentation_step3.py
"""

import os
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

sys.path.insert(0, os.path.dirname(__file__))
from make_presentation import (
    add_textbox, add_rect, add_title_bar, new_slide, make_table,
    add_bullet_box, add_callout_box,
    C_NAVY, C_IC, C_ORANGE, C_WHITE, C_BLACK, C_ALT, C_GREEN, C_RED,
    CONT_TOP, CONT_LEFT, CONT_W, CONT_H, W, H,
)
from make_presentation_step4 import make_bar_chart, add_chart_picture

PRESENTATION_ROOT = os.path.dirname(os.path.dirname(__file__))
OUT_PATH = os.path.join(PRESENTATION_ROOT, "generated", "presentation_step3_update.pptx")


def slide1_intro(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Step 3: Does Temporal Multiplexing Generalise to Mixed Ops?")

    add_callout_box(slide,
        "Design: same Plan D structure (2 shared channels, sequential sub-windows, single "
        "readout) but each query now applies a DIFFERENT boolean op, signalled by one-hot "
        "op-identity input channels. Tests whether delay-based routing survives operation "
        "heterogeneity, not just query heterogeneity.",
        CONT_LEFT, CONT_TOP, CONT_W, Inches(1.0),
        bg_color=C_ALT, border_color=C_IC, font_size=13)

    def build(tmp):
        make_bar_chart(tmp,
            ["K=1", "K=2", "K=3", "K=4"],
            {
                "w_and_d": [84.3, 74.4, 68.7, 67.6],
                "d0_control": [64.5, 60.0, 57.6, 56.7],
            },
            "Step 3: 8-Op Mixed, L1-h50 + MLP (Plan D)", "Accuracy (%)",
            threshold=90)

    add_chart_picture(slide, "step3_8op", Inches(0.5), Inches(2.1), Inches(7.0), Inches(4.1), build)

    bullets = [
        ("Delay gap persists: +19.8pp (K=1) down to +10.9pp (K=4) — same order of magnitude "
         "as same-op Plan D (+14-19% to +34%).", 0),
        ("Max K@90% = 0 for both models — but this is a DATA/TASK ceiling, not a routing "
         "failure: only ~500 samples/op (8 ops) and XOR/XNOR drag the mean down.", 0),
        ("d0 collapses toward the label prior at higher K — same structural failure mode "
         "seen in every shared-channel topology so far.", 0),
    ]
    add_bullet_box(slide, bullets, Inches(7.7), Inches(2.3), Inches(5.1), Inches(3.5), font_size=14)
    return slide


def slide2_progression(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Step 3 Progression: Data, Then Capacity, Close the Gap")

    def build_maxk(tmp):
        fig, ax = plt.subplots(figsize=(5.6, 4.0))
        fig.patch.set_facecolor("white")
        ax.set_facecolor("white")
        labels = ["4-op\n1k/op", "4-op\n16k/op\n(Dir A)", "+h=100\n(Dir C)", "+2-layer\n(Dir D)"]
        vals = [0, 2, 3, 2]
        colors = ["#C0142C", "#E8600A", "#1A7A3C", "#888888"]
        ax.bar(labels, vals, color=colors)
        for i, v in enumerate(vals):
            ax.text(i, v + 0.08, str(v), ha="center", fontsize=13, fontweight="bold")
        ax.set_ylabel("Max K @ 90% accuracy", fontsize=12)
        ax.set_title("Max K@90% Across Step 3 Directions", fontsize=13, fontweight="bold", color="#1F3864")
        ax.set_ylim(0, 3.6)
        ax.grid(True, axis="y", alpha=0.3, linestyle="--")
        for spine in ax.spines.values():
            spine.set_edgecolor("#CCCCCC")
        plt.tight_layout()
        plt.savefig(tmp, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)

    def build_k3(tmp):
        make_bar_chart(tmp,
            ["1k/op\n(Sec 17)", "16k/op\nDir A (Sec 19)", "+h=100\nDir C (Sec 20)", "+2-layer\nDir D (Sec 21)"],
            {"wad acc @ K=3": [74.9, 88.1, 90.2, 88.72]},
            "Accuracy at K=3 Across Conditions", "Accuracy (%)",
            threshold=90, colors=["#003E74"], figsize=(5.6, 4.0))

    add_chart_picture(slide, "step3_maxk", Inches(0.4), Inches(1.75), Inches(6.1), Inches(4.0), build_maxk)
    add_chart_picture(slide, "step3_k3", Inches(6.7), Inches(1.75), Inches(6.1), Inches(4.0), build_k3)

    bullets = [
        ("Direction A (4× data, 1k→16k/op): K=1 accuracy 87.7%→96.5%; Max K@90% rises 0→2 — "
         "the bottleneck was data sparsity, not the delay mechanism.", 0),
        ("Direction C (+capacity, h=50→100): Max K@90% rises 2→3, matching single-op NAND's "
         "ceiling exactly — d0 stays flat, confirming delays (not width) convert capacity into routing.", 0),
        ("Direction D (2-layer h50+h50, same 100-neuron budget): Max K@90% falls back to 2 — "
         "splitting the budget into layers hurts mixed-op routing too (same readout_in-halving "
         "confound that Section 27 later resolves for single-op NAND).", 0),
    ]
    add_bullet_box(slide, bullets, CONT_LEFT, Inches(5.85), CONT_W, Inches(1.6), font_size=13)
    return slide


def slide3_conclusions(prs):
    slide = new_slide(prs)
    add_title_bar(slide, "Step 3 Conclusions: Mixed-Op Multiplexing Confirms the Same Mechanism")

    claims = [
        ("1.  The delay advantage generalises to mixed operations",
         "+11–34pp over d0 across every Step 3 condition (8-op and 4-op, 1k and 16k samples/op) — "
         "operation heterogeneity does not break delay-based temporal routing."),
        ("2.  Data volume, not the delay mechanism, was the dominant early bottleneck",
         "4× more training data (1k→16k/op) alone lifted Max K@90% from 0 to 2 — w_and_d absorbs "
         "extra data; d0_control stays flat, confirming delays are what convert data into capability."),
        ("3.  Width helps, naive depth hurts — same lesson as single-op NAND, now confirmed for mixed ops",
         "h=50→100 (Dir C) raises Max K@90% to 3, matching single-op NAND. Splitting into 2 layers at "
         "the same 100-neuron budget (Dir D) drops it back to 2 — depth only helps once the readout-input "
         "confound is fixed (see Section 27, single-op, later in this talk)."),
    ]
    colors = [
        (RGBColor(0xE8, 0xF0, 0xFE), C_IC),
        (RGBColor(0xFF, 0xF3, 0xE0), C_ORANGE),
        (RGBColor(0xD6, 0xE8, 0xD6), C_GREEN),
    ]
    top_start = CONT_TOP
    box_h = Inches(1.6)
    gap = Inches(0.15)
    for idx, (title, body) in enumerate(claims):
        bg, border = colors[idx]
        add_callout_box(slide, f"▶  {title}\n    {body}",
                        CONT_LEFT, top_start + idx * (box_h + gap),
                        CONT_W, box_h, bg_color=bg, border_color=border, font_size=13)

    add_textbox(slide,
                "This sets up Step 4: having shown the mechanism generalises across operations, "
                "Step 4 next asks whether it generalises across input/output STRUCTURE.",
                CONT_LEFT, Inches(6.55), CONT_W, Inches(0.6),
                font_size=14, bold=True, color=C_NAVY)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    print("Building Step 3 slides...")
    slide1_intro(prs); print("  [1/3] Intro + 8-op results")
    slide2_progression(prs); print("  [2/3] Progression (Dir A/C/D)")
    slide3_conclusions(prs); print("  [3/3] Conclusions")

    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}")


if __name__ == "__main__":
    main()
